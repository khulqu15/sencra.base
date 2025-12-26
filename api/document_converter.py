from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
import tempfile, os, csv, subprocess
import tempfile
import csv
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import xlrd, xlwt
from pdf2image import convert_from_path
import xlrd
import xlwt
from utils import save_upload_to_temp, extract_pdf_tables, libreoffice_convert, doc_to_docx
from openpyxl import Workbook, load_workbook
import aiofiles
import csv
from io import StringIO

router = APIRouter(prefix="/convert", tags=["Document Converter"])

EXT_PDF  = ".pdf"
EXT_DOCX = ".docx"
EXT_DOC  = ".doc"
EXT_XLSX = ".xlsx"
EXT_XLS  = ".xls"
EXT_CSV  = ".csv"
EXT_PPTX = ".pptx"

MIME_PDF  = "application/pdf"
MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MIME_DOC  = "application/msword"
MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
MIME_XLS  = "application/vnd.ms-excel"
MIME_CSV  = "text/csv"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
EXCEPTION_NO_TABLE = "No tables found in DOCX"

@router.post("/csv-to-xlsx")
async def convert_csv_to_xlsx(file: UploadFile = File(...)):
    csv_file = save_upload_to_temp(file, EXT_CSV)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_XLSX).name
    try:
        wb = Workbook()
        ws = wb.active
        async with aiofiles.open(csv_file, newline="", encoding="utf-8") as f:
            content = await f.read()
            
        reader = csv.reader(StringIO(content))
        for row in reader:
            ws.append(row)
        wb.save(out)
        return FileResponse(out, MIME_XLSX, file.filename.replace(EXT_CSV, EXT_XLSX))
    finally:
        os.unlink(csv_file)

@router.post("/csv-to-docx")
async def convert_csv_to_docx(file: UploadFile = File(...)):
    csv_file = save_upload_to_temp(file, EXT_CSV)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    try:
        async with aiofiles.open(csv_file, "r", encoding="utf-8") as f:
            content = await f.read()
        rows = list(csv.reader(StringIO(content)))
        doc = Document()
        table = doc.add_table(rows=len(rows), cols=len(rows[0]))
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                table.rows[r].cells[c].text = val
        doc.save(out)
        return FileResponse(out, MIME_DOCX, file.filename.replace(EXT_CSV, EXT_DOCX))
    finally:
        os.unlink(csv_file)

@router.post("/csv-to-doc")
async def convert_csv_to_doc(file: UploadFile = File(...)):
    csv_file = save_upload_to_temp(file, EXT_CSV)
    tmp_docx = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    out_dir = tempfile.mkdtemp()
    out_doc = os.path.join(out_dir, file.filename.replace(EXT_CSV, ".doc"))
    try:
        async with aiofiles.open(csv_file, "r", encoding="utf-8") as f:
            content = await f.read()
        rows = list(csv.reader(StringIO(content)))
        doc = Document()
        table = doc.add_table(rows=len(rows), cols=len(rows[0]))
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                table.rows[r].cells[c].text = val
        doc.save(tmp_docx)
        libreoffice_convert(tmp_docx, out_dir, "doc")
        return FileResponse(out_doc, MIME_DOC, file.filename.replace(EXT_CSV, ".doc"))
    finally:
        os.unlink(csv_file)
        if os.path.exists(tmp_docx):
            os.unlink(tmp_docx)

@router.post("/csv-to-pdf")
async def convert_csv_to_pdf(file: UploadFile = File(...)):
    csv_file = save_upload_to_temp(file, EXT_CSV)
    tmp_xlsx = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_XLSX).name
    out_dir = tempfile.mkdtemp()
    out_pdf = os.path.join(out_dir, file.filename.replace(EXT_CSV, ".pdf"))
    try:
        wb = Workbook()
        ws = wb.active
        async with aiofiles.open(csv_file, newline="", encoding="utf-8") as f:
            content = await f.read()
        reader = csv.reader(StringIO(content))
        for row in reader:
            ws.append(row)
        wb.save(tmp_xlsx)
        libreoffice_convert(tmp_xlsx, out_dir, "pdf")
        return FileResponse(out_pdf, MIME_PDF, file.filename.replace(EXT_CSV, ".pdf"))
    finally:
        os.unlink(csv_file)

@router.post("/csv-to-pptx")
async def convert_csv_to_pptx(file: UploadFile = File(...)):
    csv_file = save_upload_to_temp(file, EXT_CSV)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_PPTX).name
    try:
        async with aiofiles.open(csv_file, "r", encoding="utf-8") as f:
            content = await f.read()
        rows = list(csv.reader(StringIO(content)))
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = slide.shapes.add_table(
            rows=len(rows),
            cols=len(rows[0]),
            left=Inches(0.5),
            top=Inches(1),
            width=prs.slide_width - Inches(1),
            height=prs.slide_height - Inches(2),
        ).table
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                table.cell(r, c).text = val

        prs.save(out)
        return FileResponse(out, MIME_PPTX, file.filename.replace(EXT_CSV, EXT_PPTX))
    finally:
        os.unlink(csv_file)

## |-------------------------|
## | XLS CONVERTION FUNCTION |
## |-------------------------|

@router.post("/xls-to-xlsx")
async def convert_xls_to_xlsx(file: UploadFile = File(...)):
    xls = save_upload_to_temp(file, ".xls")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_XLSX).name
    try:
        book = xlrd.open_workbook(xls)
        wb = Workbook()
        ws = wb.active
        sheet = book.sheet_by_index(0)
        for r in range(sheet.nrows):
            ws.append(sheet.row_values(r))
        wb.save(out)
        return FileResponse(out, MIME_XLSX, file.filename.replace(".xls", EXT_XLSX))
    finally:
        os.unlink(xls)

@router.post("/xls-to-csv")
async def convert_xls_to_csv(file: UploadFile = File(...)):
    xls = save_upload_to_temp(file, ".xls")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_CSV).name
    try:
        book = xlrd.open_workbook(xls)
        sheet = book.sheet_by_index(0)
        buffer = StringIO()
        writer = csv.writer(buffer)
        for r in range(sheet.nrows):
            writer.writerow(sheet.row_values(r))

        async with aiofiles.open(out, "w", encoding="utf-8", newline="") as f:
            await f.write(buffer.getvalue())
        return FileResponse(out, MIME_CSV, file.filename.replace(".xls", EXT_CSV))
    finally:
        os.unlink(xls)

@router.post("/xls-to-docx")
async def convert_xls_to_docx(file: UploadFile = File(...)):
    xls = save_upload_to_temp(file, ".xls")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    try:
        book = xlrd.open_workbook(xls)
        sheet = book.sheet_by_index(0)
        doc = Document()
        table = doc.add_table(rows=sheet.nrows, cols=sheet.ncols)
        for r in range(sheet.nrows):
            for c in range(sheet.ncols):
                table.rows[r].cells[c].text = str(sheet.cell_value(r, c))
        doc.save(out)
        return FileResponse(out, MIME_DOCX, file.filename.replace(".xls", EXT_DOCX))
    finally:
        os.unlink(xls)

@router.post("/xls-to-doc")
async def convert_xls_to_doc(file: UploadFile = File(...)):
    xls = save_upload_to_temp(file, ".xls")
    tmp_docx = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    out_dir = tempfile.mkdtemp()
    out_doc = os.path.join(out_dir, file.filename.replace(".xls", ".doc"))
    try:
        book = xlrd.open_workbook(xls)
        sheet = book.sheet_by_index(0)
        doc = Document()
        table = doc.add_table(rows=sheet.nrows, cols=sheet.ncols)
        for r in range(sheet.nrows):
            for c in range(sheet.ncols):
                table.rows[r].cells[c].text = str(sheet.cell_value(r, c))
        doc.save(tmp_docx)
        libreoffice_convert(tmp_docx, out_dir, "doc")
        return FileResponse(out_doc, MIME_DOC, file.filename.replace(".xls", ".doc"))
    finally:
        os.unlink(xls)
        if os.path.exists(tmp_docx):
            os.unlink(tmp_docx)

@router.post("/xls-to-pdf")
async def convert_xls_to_pdf(file: UploadFile = File(...)):
    xls = save_upload_to_temp(file, ".xls")
    out_dir = tempfile.mkdtemp()
    out_pdf = os.path.join(out_dir, file.filename.replace(".xls", ".pdf"))
    try:
        libreoffice_convert(xls, out_dir, "pdf")
        return FileResponse(out_pdf, MIME_PDF, file.filename.replace(".xls", ".pdf"))
    finally:
        os.unlink(xls)

@router.post("/xls-to-pptx")
async def convert_xls_to_pptx(file: UploadFile = File(...)):
    xls = save_upload_to_temp(file, ".xls")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_PPTX).name
    try:
        book = xlrd.open_workbook(xls)
        sheet = book.sheet_by_index(0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = slide.shapes.add_table(
            sheet.nrows, sheet.ncols,
            Inches(0.5), Inches(1),
            prs.slide_width - Inches(1),
            prs.slide_height - Inches(2)
        ).table
        for r in range(sheet.nrows):
            for c in range(sheet.ncols):
                table.cell(r, c).text = str(sheet.cell_value(r, c))
        prs.save(out)
        return FileResponse(out, MIME_PPTX, file.filename.replace(".xls", EXT_PPTX))
    finally:
        os.unlink(xls)


## |-------------------------|
## | DOC CONVERTION FUNCTION |
## |-------------------------|

@router.post("/doc-to-pdf")
async def convert_doc_to_pdf(file: UploadFile = File(...)):
    doc = save_upload_to_temp(file, ".doc")
    out_dir = tempfile.mkdtemp()
    out_pdf = os.path.join(out_dir, file.filename.replace(".doc", ".pdf"))
    try:
        libreoffice_convert(doc, out_dir, "pdf")
        return FileResponse(out_pdf, MIME_PDF, file.filename.replace(".doc", ".pdf"))
    finally:
        os.unlink(doc)

@router.post("/doc-to-docx")
async def convert_doc_to_docx(file: UploadFile = File(...)):
    doc = save_upload_to_temp(file, ".doc")
    out_dir = tempfile.mkdtemp()
    out_docx = os.path.join(out_dir, file.filename.replace(".doc", EXT_DOCX))
    try:
        libreoffice_convert(doc, out_dir, "docx")
        return FileResponse(out_docx, MIME_DOCX, file.filename.replace(".doc", EXT_DOCX))
    finally:
        os.unlink(doc)

@router.post("/doc-to-xlsx")
async def convert_doc_to_xlsx(file: UploadFile = File(...)):
    doc = save_upload_to_temp(file, ".doc")
    docx = doc_to_docx(doc)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_XLSX).name
    try:
        document = Document(docx)
        wb = Workbook()
        ws = wb.active
        for table in document.tables:
            for row in table.rows:
                ws.append([cell.text for cell in row.cells])
            ws.append([])
        wb.save(out)
        return FileResponse(out, MIME_XLSX, file.filename.replace(".doc", EXT_XLSX))
    finally:
        os.unlink(doc)

@router.post("/doc-to-xls")
async def convert_doc_to_xls(file: UploadFile = File(...)):
    doc = save_upload_to_temp(file, ".doc")
    docx = doc_to_docx(doc)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".xls").name
    try:
        document = Document(docx)
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Tables")
        r = 0
        for table in document.tables:
            for row in table.rows:
                for c, cell in enumerate(row.cells):
                    ws.write(r, c, cell.text)
                r += 1
            r += 1
        wb.save(out)
        return FileResponse(out, MIME_XLS, file.filename.replace(".doc", ".xls"))
    finally:
        os.unlink(doc)

@router.post("/doc-to-csv")
async def convert_doc_to_csv(file: UploadFile = File(...)):
    doc = save_upload_to_temp(file, ".doc")
    docx = doc_to_docx(doc)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_CSV).name
    try:
        document = Document(docx)
        buffer = StringIO()
        writer = csv.writer(buffer)
        for table in document.tables:
            for row in table.rows:
                writer.writerow([cell.text for cell in row.cells])
            writer.writerow([])
        async with aiofiles.open(out, "w", encoding="utf-8", newline="") as f:
            await f.write(buffer.getvalue())
        return FileResponse(out, MIME_CSV, file.filename.replace(".doc", EXT_CSV))
    finally:
        os.unlink(doc)

@router.post("/doc-to-pptx")
async def convert_doc_to_pptx(file: UploadFile = File(...)):
    doc = save_upload_to_temp(file, ".doc")
    docx = doc_to_docx(doc)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_PPTX).name
    try:
        document = Document(docx)
        prs = Presentation()
        for para in document.paragraphs:
            if para.text.strip():
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = para.text
        prs.save(out)
        return FileResponse(out, MIME_PPTX, file.filename.replace(".doc", EXT_PPTX))
    finally:
        os.unlink(doc)


## |--------------------------|
## | PPTX CONVERTION FUNCTION |
## |--------------------------|

@router.post("/pptx-to-pdf")
async def convert_pptx_to_pdf(file: UploadFile = File(...)):
    pptx = save_upload_to_temp(file, EXT_PPTX)
    out_dir = tempfile.mkdtemp()
    out_pdf = os.path.join(out_dir, file.filename.replace(EXT_PPTX, ".pdf"))
    try:
        libreoffice_convert(pptx, out_dir, "pdf")
        return FileResponse( out_pdf, media_type=MIME_PDF, filename=file.filename.replace(EXT_PPTX, ".pdf"))
    finally:
        os.unlink(pptx)

@router.post("/pptx-to-xls")
async def convert_pptx_to_xls(file: UploadFile = File(...)):
    pptx = save_upload_to_temp(file, EXT_PPTX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".xls").name
    try:
        prs = Presentation(pptx)
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Slides")
        ws.write(0, 0, "Slide")
        ws.write(0, 1, "Text")
        row_idx = 1
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            ws.write(row_idx, 0, i + 1)
            ws.write(row_idx, 1, " | ".join(texts))
            row_idx += 1
        wb.save(out)
        return FileResponse(out,media_type=MIME_XLS,filename=file.filename.replace(EXT_PPTX, ".xls"))
    finally:
        os.unlink(pptx)

@router.post("/pptx-to-docx")
async def convert_pptx_to_docx(file: UploadFile = File(...)):
    pptx = save_upload_to_temp(file, EXT_PPTX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    try:
        prs = Presentation(pptx)
        doc = Document()
        for i, slide in enumerate(prs.slides):
            doc.add_heading(f"Slide {i+1}", level=2)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    doc.add_paragraph(shape.text)
            doc.add_page_break()
        doc.save(out)
        return FileResponse(out, media_type=MIME_DOCX, filename=file.filename.replace(EXT_PPTX, EXT_DOCX))
    finally:
        os.unlink(pptx)

@router.post("/pptx-to-doc")
async def convert_pptx_to_doc(file: UploadFile = File(...)):
    pptx = save_upload_to_temp(file, EXT_PPTX)
    tmp_docx = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    out_dir = tempfile.mkdtemp()
    out_doc = os.path.join(out_dir, file.filename.replace(EXT_PPTX, ".doc"))
    try:
        prs = Presentation(pptx)
        doc = Document()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    doc.add_paragraph(shape.text)
            doc.add_page_break()
        doc.save(tmp_docx)
        libreoffice_convert(tmp_docx, out_dir, "doc")
        return FileResponse(out_doc, media_type=MIME_DOC, filename=file.filename.replace(EXT_PPTX, ".doc"))
    finally:
        os.unlink(pptx)
        if os.path.exists(tmp_docx):
            os.unlink(tmp_docx)

@router.post("/pptx-to-xlsx")
async def convert_pptx_to_xlsx(file: UploadFile = File(...)):
    pptx = save_upload_to_temp(file, EXT_PPTX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_XLSX).name
    try:
        prs = Presentation(pptx)
        wb = Workbook()
        ws = wb.active
        ws.append(["Slide", "Text"])
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            ws.append([i + 1, " | ".join(texts)])
        wb.save(out)
        return FileResponse(out, media_type=MIME_XLSX, filename=file.filename.replace(EXT_PPTX, EXT_XLSX))
    finally:
        os.unlink(pptx)
        
@router.post("/pptx-to-csv")
async def convert_pptx_to_csv(file: UploadFile = File(...)):
    pptx = save_upload_to_temp(file, EXT_PPTX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_CSV).name
    try:
        prs = Presentation(pptx)
        buffer = StringIO()
        writer = csv.writer(buffer)
        writer.writerow(["Slide", "Text"])
        for i, slide in enumerate(prs.slides):
            texts = [
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            writer.writerow([i + 1, " | ".join(texts)])
        async with aiofiles.open(out, "w", encoding="utf-8", newline="") as f:
            await f.write(buffer.getvalue())
        return FileResponse(out, media_type=MIME_CSV, filename=file.filename.replace(EXT_PPTX, EXT_CSV))
    finally:
        os.unlink(pptx)
        

## |--------------------------|
## | XLSX CONVERTION FUNCTION |
## |--------------------------|

@router.post("/xlsx-to-pdf")
async def convert_xlsx_to_pdf(file: UploadFile = File(...)):
    xlsx = save_upload_to_temp(file, EXT_XLSX)
    out_dir = tempfile.mkdtemp()
    out_pdf = os.path.join(out_dir, file.filename.replace(EXT_XLSX, ".pdf"))
    try:
        libreoffice_convert(xlsx, out_dir, "pdf")
        return FileResponse(out_pdf, media_type=MIME_PDF, filename=file.filename.replace(EXT_XLSX, ".pdf"))
    finally:
        os.unlink(xlsx)

@router.post("/xlsx-to-docx")
async def convert_xlsx_to_docx(file: UploadFile = File(...)):
    xlsx = save_upload_to_temp(file, EXT_XLSX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    try:
        wb = load_workbook(xlsx, data_only=True)
        doc = Document()
        for sheet in wb.worksheets:
            doc.add_heading(sheet.title, level=1)
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join("" if c is None else str(c) for c in row)
                doc.add_paragraph(line)
            doc.add_page_break()
        doc.save(out)
        return FileResponse(out, media_type=MIME_DOCX, filename=file.filename.replace(EXT_XLSX, EXT_DOCX))
    finally:
        os.unlink(xlsx)

@router.post("/xlsx-to-doc")
async def convert_xlsx_to_doc(file: UploadFile = File(...)):
    xlsx = save_upload_to_temp(file, EXT_XLSX)
    tmp_docx = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    out_dir = tempfile.mkdtemp()
    out_doc = os.path.join(out_dir, file.filename.replace(EXT_XLSX, ".doc"))
    try:
        wb = load_workbook(xlsx, data_only=True)
        doc = Document()
        for sheet in wb.worksheets:
            doc.add_heading(sheet.title, level=1)
            for row in sheet.iter_rows(values_only=True):
                doc.add_paragraph("\t".join("" if c is None else str(c) for c in row))
            doc.add_page_break()
        doc.save(tmp_docx)
        libreoffice_convert(tmp_docx, out_dir, "doc")
        return FileResponse(out_doc, media_type=MIME_DOC, filename=file.filename.replace(EXT_XLSX, ".doc"))
    finally:
        os.unlink(xlsx)
        if os.path.exists(tmp_docx):
            os.unlink(tmp_docx)

@router.post("/xlsx-to-csv")
async def convert_xlsx_to_csv(file: UploadFile = File(...)):
    xlsx = save_upload_to_temp(file, EXT_XLSX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_CSV).name
    try:
        wb = load_workbook(xlsx, data_only=True)
        sheet = wb.active
        buffer = StringIO()
        writer = csv.writer(buffer)
        for row in sheet.iter_rows(values_only=True):
            writer.writerow(row)
        async with aiofiles.open(out, "w", encoding="utf-8", newline="") as f:
            await f.write(buffer.getvalue())
        return FileResponse(out, media_type=MIME_CSV, filename=file.filename.replace(EXT_XLSX, EXT_CSV))
    finally:
        os.unlink(xlsx)
        
@router.post("/xlsx-to-pptx")
async def convert_xlsx_to_pptx(file: UploadFile = File(...)):
    xlsx = save_upload_to_temp(file, EXT_XLSX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_PPTX).name
    try:
        wb = load_workbook(xlsx, data_only=True)
        prs = Presentation()
        for sheet in wb.worksheets:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sheet.title
            content = "\n".join(
                "\t".join("" if c is None else str(c) for c in row)
                for row in sheet.iter_rows(values_only=True)
            )
            slide.placeholders[1].text = content[:5000]

        prs.save(out)
        return FileResponse(out, media_type=MIME_PPTX, filename=file.filename.replace(EXT_XLSX, EXT_PPTX))
    finally:
        os.unlink(xlsx)

@router.post("/xlsx-to-xls")
async def convert_xlsx_to_xls(file: UploadFile = File(...)):
    xlsx = save_upload_to_temp(file, EXT_XLSX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".xls").name
    try:
        wb_in = load_workbook(xlsx, data_only=True)
        wb_out = xlwt.Workbook()
        ws_out = wb_out.add_sheet("Sheet1")
        r = 0
        for row in wb_in.active.iter_rows(values_only=True):
            for c, val in enumerate(row):
                ws_out.write(r, c, "" if val is None else str(val))
            r += 1
        wb_out.save(out)
        return FileResponse(out, media_type=MIME_XLS, filename=file.filename.replace(EXT_XLSX, ".xls"))
    finally:
        os.unlink(xlsx)

        
## |--------------------------|
## | DOCX CONVERTION FUNCTION |
## |--------------------------|

@router.post("/docx-to-xlsx")
async def convert_docx_to_xlsx(file: UploadFile = File(...)):
    docx = save_upload_to_temp(file, EXT_DOCX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_XLSX).name
    try:
        doc = Document(docx)
        wb = Workbook()
        ws = wb.active
        found = False
        for table in doc.tables:
            found = True
            for row in table.rows:
                ws.append([cell.text for cell in row.cells])
            ws.append([])
        if not found:
            raise HTTPException(status_code=400, detail=EXCEPTION_NO_TABLE)
        wb.save(out)
        return FileResponse(out, media_type=MIME_XLSX, filename=file.filename.replace(EXT_DOCX, EXT_XLSX))
    finally:
        os.unlink(docx)

@router.post("/docx-to-pptx")
async def convert_docx_to_pptx(file: UploadFile = File(...)):
    docx = save_upload_to_temp(file, EXT_DOCX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_PPTX).name
    try:
        doc = Document(docx)
        prs = Presentation()
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = para.text
        prs.save(out)
        return FileResponse(out, media_type=MIME_PPTX, filename=file.filename.replace(EXT_DOCX, EXT_PPTX))
    finally:
        os.unlink(docx)
        
@router.post("/docx-to-pdf")
async def convert_docx_to_pdf(file: UploadFile = File(...)):
    docx = save_upload_to_temp(file, EXT_DOCX)
    out_dir = tempfile.mkdtemp()
    out_pdf = os.path.join(out_dir, file.filename.replace(EXT_DOCX, ".pdf"))
    try:
        libreoffice_convert(docx, out_dir, "pdf")
        return FileResponse(out_pdf, media_type=MIME_PDF, filename=file.filename.replace(EXT_DOCX, ".pdf"))
    finally:
        os.unlink(docx)

@router.post("/docx-to-doc")
async def convert_docx_to_doc(file: UploadFile = File(...)):
    docx = save_upload_to_temp(file, EXT_DOCX)
    out_dir = tempfile.mkdtemp()
    out_doc = os.path.join(out_dir, file.filename.replace(EXT_DOCX, ".doc"))
    try:
        libreoffice_convert(docx, out_dir, "doc")
        return FileResponse(out_doc, media_type=MIME_DOC, filename=file.filename.replace(EXT_DOCX, ".doc"))
    finally:
        os.unlink(docx)

@router.post("/docx-to-xls")
async def convert_docx_to_xls(file: UploadFile = File(...)):
    docx = save_upload_to_temp(file, EXT_DOCX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".xls").name

    try:
        doc = Document(docx)
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Tables")
        r = 0
        found = False
        for table in doc.tables:
            found = True
            for row in table.rows:
                for c, cell in enumerate(row.cells):
                    ws.write(r, c, cell.text)
                r += 1
            r += 1
        if not found:
            raise HTTPException(status_code=400, detail=EXCEPTION_NO_TABLE)
        wb.save(out)
        return FileResponse(out, media_type=MIME_XLS, filename=file.filename.replace(EXT_DOCX, ".xls"))
    finally:
        os.unlink(docx)

@router.post("/docx-to-csv")
async def convert_docx_to_csv(file: UploadFile = File(...)):
    docx = save_upload_to_temp(file, EXT_DOCX)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_CSV).name
    try:
        doc = Document(docx)
        buffer = StringIO()
        writer = csv.writer(buffer)
        for table in doc.tables:
            for row in table.rows:
                writer.writerow([cell.text for cell in row.cells])
            writer.writerow([])
        async with aiofiles.open(out, "w", encoding="utf-8", newline="") as f:
            await f.write(buffer.getvalue())
        return FileResponse(out, media_type=MIME_CSV, filename=file.filename.replace(EXT_DOCX, EXT_CSV))
    finally:
        os.unlink(docx)

        
## |-------------------------|
## | PDF CONVERTION FUNCTION |
## |-------------------------|


@router.post("/pdf-to-xlsx")
async def convert_pdf_to_xlsx(file: UploadFile = File(...)):
    pdf = save_upload_to_temp(file, ".pdf")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    try:
        tables = extract_pdf_tables(pdf)
        wb = Workbook()
        ws = wb.active
        for table in tables:
            for row in table.df.values.tolist():
                ws.append(row)
            ws.append([])
        wb.save(out.name)
        return FileResponse(out, media_type=MIME_XLSX, filename=file.filename.replace(".pdf", EXT_XLSX))
    finally:
        os.unlink(pdf)


@router.post("/pdf-to-docx")
async def convert_pdf_to_docx(file: UploadFile = File(...)):
    pdf = save_upload_to_temp(file, ".pdf")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    try:
        tables = extract_pdf_tables(pdf)
        doc = Document()
        for i, table in enumerate(tables):
            doc.add_heading(f"Table {i + 1}", level=2)
            rows, cols = table.df.shape
            doc_table = doc.add_table(rows=rows, cols=cols)
            for r in range(rows):
                for c in range(cols):
                    doc_table.rows[r].cells[c].text = str(table.df.iat[r, c])
            doc.add_paragraph() 
        doc.save(out)
        return FileResponse(out,media_type=MIME_DOCX,filename=file.filename.replace(".pdf", EXT_DOCX))
    finally:
        os.unlink(pdf)
        

@router.post("/pdf-to-doc")
async def convert_pdf_to_doc(file: UploadFile = File(...)):
    pdf = save_upload_to_temp(file, ".pdf")
    tmp_docx = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_DOCX).name
    out_doc = tempfile.NamedTemporaryFile(delete=False, suffix=".doc").name
    try:
        tables = extract_pdf_tables(pdf)
        doc = Document()
        for i, table in enumerate(tables):
            doc.add_heading(f"Table {i + 1}", level=2)
            rows, cols = table.df.shape
            doc_table = doc.add_table(rows, cols)
            for r in range(rows):
                for c in range(cols):
                    doc_table.rows[r].cells[c].text = str(table.df.iat[r, c])
            doc.add_paragraph()
        doc.save(tmp_docx)
        libreoffice_convert(tmp_docx, os.path.dirname(out_doc), "doc")
        return FileResponse(out_doc, media_type=MIME_DOC, filename=file.filename.replace(".pdf", ".doc"))
    finally:
        os.unlink(pdf)
        if os.path.exists(tmp_docx):
            os.unlink(tmp_docx)

@router.post("/pdf-to-csv")
async def convert_pdf_to_csv(file: UploadFile = File(...)):
    pdf = save_upload_to_temp(file, ".pdf")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=EXT_CSV).name
    try:
        tables = extract_pdf_tables(pdf)
        buffer = StringIO()
        writer = csv.writer(buffer)
        for i, table in enumerate(tables):
            for row in table.df.values.tolist():
                writer.writerow(row)
            if i < tables.n - 1:
                writer.writerow([])
        async with aiofiles.open(out, "w", encoding="utf-8", newline="") as f:
            await f.write(buffer.getvalue())
        return FileResponse(out.name, media_type=MIME_CSV, filename=file.filename.replace(".pdf", EXT_CSV))
    finally:
        os.unlink(pdf)

@router.post("/pdf-to-xls")
async def convert_pdf_to_xls(file: UploadFile = File(...)):
    pdf = save_upload_to_temp(file, ".pdf")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".xls").name
    try:
        tables = extract_pdf_tables(pdf)
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Tables")
        r = 0
        for table in tables:
            for row in table.df.values.tolist():
                for col_idx, value in enumerate(row):
                    ws.write(r, col_idx, value)
                r += 1
            r += 1 
        wb.save(out)
        return FileResponse(out, media_type=MIME_XLS, filename=file.filename.replace(".pdf", ".xls"),)
    finally:
        os.unlink(pdf)
        
@router.post("/pdf-to-pptx")
async def convert_pdf_to_pptx(file: UploadFile = File(...)):
    pdf = save_upload_to_temp(file, ".pdf")
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".xls").name
    try:
        images = convert_from_path(pdf)
        prs = Presentation()
        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            img_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            img.save(img_path, "PNG")
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.unlink(img_path)
        prs.save(out.name)
        return FileResponse(out, media_type=MIME_PPTX, filename=file.filename.replace(".pdf", EXT_PPTX))
    finally:
        os.unlink(pdf)