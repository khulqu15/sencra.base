from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import FileResponse
from pdf2docx import Converter
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pdf2image import convert_from_path
import tempfile
import os
import subprocess

def pptx_to_pdf_linux(input_path: str, output_path: str):
    try:
        subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", os.path.dirname(output_path),
            input_path
        ], check=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Failed to convert PPTX to PDF: {e}")
    
app = FastAPI()

@app.post("/convert/pdf-to-docx")
async def convert_pdf_to_docx(file: UploadFile = File(...)):
    if not file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="File must be PDF")

    input_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    input_temp.write(await file.read())
    input_temp.close()
    output_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    output_temp.close()

    try:
        cv = Converter(input_temp.name)
        cv.convert(output_temp.name, start=0, end=None)
        cv.close()
        return FileResponse(output_temp.name, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename=file.filename.replace(".pdf", ".docx"))

    finally:
        os.unlink(input_temp.name)


@app.post("/convert/docx-to-xlsx")
async def convert_docx_to_xlsx(file: UploadFile = File(...)):
    if not file.filename.endswith(".docx"):
        raise HTTPException(status_code=400, detail="File must be DOCX")

    input_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    input_temp.write(await file.read())
    input_temp.close()

    output_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")
    output_temp.close()

    try:
        doc = Document(input_temp.name)

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"

        table_count = 0
        for table in doc.tables:
            table_count += 1
            for row in table.rows:
                ws.append([cell.text for cell in row.cells])
            ws.append([]) 

        if table_count == 0:
            raise HTTPException(status_code=400, detail="No tables found in DOCX")
        wb.save(output_temp.name)
        return FileResponse(
            output_temp.name,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=file.filename.replace(".docx", ".xlsx")
        )

    finally:
        os.unlink(input_temp.name)
        
@app.post("/convert/xlsx-to-docx")
async def convert_xlsx_to_docx(file: UploadFile = File(...)):
    if not file.filename.endswith((".xlsx", ".xls")):
        raise HTTPException(status_code=400, detail="File must be XLSX or XLS")

    input_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")
    input_temp.write(await file.read())
    input_temp.close()

    output_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    output_temp.close()

    try:
        wb = load_workbook(input_temp.name, data_only=True)
        doc = Document()

        for sheet in wb.worksheets:
            doc.add_heading(sheet.title, level=1)
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join([str(cell) if cell is not None else "" for cell in row])
                doc.add_paragraph(line)
            doc.add_page_break() 

        doc.save(output_temp.name)

        return FileResponse(
            output_temp.name,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=file.filename.rsplit(".", 1)[0] + ".docx"
        )

    finally:
        os.unlink(input_temp.name)
        
# --- PDF → PPTX ---
@app.post("/convert/pdf-to-pptx")
async def convert_pdf_to_pptx(file: UploadFile = File(...)):
    if not file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="File must be PDF")

    input_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    input_temp.write(await file.read())
    input_temp.close()
    output_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    output_temp.close()

    try:
        images = convert_from_path(input_temp.name)
        prs = Presentation()
        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            img_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            img.save(img_path, "PNG")
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.unlink(img_path)
        prs.save(output_temp.name)
        return FileResponse(output_temp.name, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=file.filename.replace(".pdf", ".pptx"))
    finally:
        os.unlink(input_temp.name)

# --- DOCX → PPTX ---
@app.post("/convert/docx-to-pptx")
async def convert_docx_to_pptx(file: UploadFile = File(...)):
    if not file.filename.endswith(".docx"):
        raise HTTPException(status_code=400, detail="File must be DOCX")

    input_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    input_temp.write(await file.read())
    input_temp.close()
    output_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    output_temp.close()

    try:
        doc = Document(input_temp.name)
        prs = Presentation()
        for para in doc.paragraphs:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = para.text
        prs.save(output_temp.name)
        return FileResponse(output_temp.name, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=file.filename.replace(".docx", ".pptx"))
    finally:
        os.unlink(input_temp.name)

@app.post("/convert/pptx-to-pdf")
async def convert_pptx_to_pdf(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="File must be PPTX")

    input_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    input_temp.write(await file.read())
    input_temp.close()

    output_dir = tempfile.mkdtemp()
    output_path = os.path.join(output_dir, file.filename.replace(".pptx", ".pdf"))

    try:
        pptx_to_pdf_linux(input_temp.name, output_path)
        return FileResponse(
            output_path,
            media_type="application/pdf",
            filename=file.filename.replace(".pptx", ".pdf")
        )
    finally:
        os.unlink(input_temp.name)
        if os.path.exists(output_path):
            os.unlink(output_path)
        os.rmdir(output_dir)


@app.post("/convert/pptx-to-docx")
async def convert_pptx_to_docx(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="File must be PPTX")

    input_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    input_temp.write(await file.read())
    input_temp.close()
    output_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    output_temp.close()

    try:
        prs = Presentation(input_temp.name)
        doc = Document()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    doc.add_paragraph(shape.text)
            doc.add_page_break()
        doc.save(output_temp.name)
        return FileResponse(output_temp.name, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename=file.filename.replace(".pptx", ".docx"))
    finally:
        os.unlink(input_temp.name)